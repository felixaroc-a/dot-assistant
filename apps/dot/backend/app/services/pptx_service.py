"""Servicio de lectura y generación de presentaciones PowerPoint (.pptx)."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from app.services.document_output_service import (
    BRAND_DARK,
    resolve_output_path,
    sanitize_document_title,
)

log = logging.getLogger("dot.pptx")


def _get_dot_work_dir() -> Path:
    from app.services.document_output_service import get_desktop_work_dir

    return get_desktop_work_dir()


def _sanitize_filename(name: str) -> str:
    return sanitize_document_title(name, fallback="presentacion-dot")


def _add_chart(
    slide: Any,
    chart_data: dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type_map: dict[str, Any] = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart_type = chart_data.get("type", "bar")
    chart_enum = chart_type_map.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)

    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data.get("categories", [])

    for series in chart_data.get("series", []):
        chart_data_obj.add_series(series.get("name", ""), series.get("values", []))

    chart_frame = slide.shapes.add_chart(chart_enum, left, top, width, height, chart_data_obj)
    chart = chart_frame.chart
    chart.has_legend = True

    chart_title = chart_data.get("title")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title


def _add_title_slide(prs: Any, title: str) -> None:
    """Slide de portada con título grande y subtítulo DOT."""
    from datetime import datetime

    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.size = Pt(36)
            para.font.bold = True
            para.font.color.rgb = RGBColor.from_string(BRAND_DARK)
            para.alignment = PP_ALIGN.CENTER

    subtitle_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            subtitle_shape = shape
            break

    subtitle_text = f"Generado por DOT — {datetime.now().strftime('%d/%m/%Y')}"
    if subtitle_shape is not None:
        tf = subtitle_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor.from_string("666666")
        p.alignment = PP_ALIGN.CENTER
    else:
        box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER


def _build_default_template(prs: Any, slides_data: list[dict[str, Any]]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide_layouts = prs.slide_layouts
    content_layout = slide_layouts[1] if len(slide_layouts) > 1 else slide_layouts[0]

    for sd in slides_data:
        slide = prs.slides.add_slide(content_layout)

        title = str(sd.get("title", "")).strip()
        content = str(sd.get("content", "")).strip()

        if title and slide.shapes.title:
            slide.shapes.title.text = title
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(28)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor.from_string(BRAND_DARK)
                paragraph.alignment = PP_ALIGN.LEFT

        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break

        if body_shape is not None and content:
            tf = body_shape.text_frame
            tf.clear()
            for line_idx, line in enumerate(content.split("\n")):
                stripped = line.strip()
                if not stripped:
                    continue
                p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                if stripped.startswith("- ") or stripped.startswith("* "):
                    p.text = stripped[2:]
                    p.level = 0
                else:
                    p.text = stripped
                    p.level = 0
                p.font.size = Pt(18)
                p.font.name = "Calibri"

        image_urls = sd.get("image_urls", [])
        if image_urls:
            for i, img_path in enumerate(image_urls):
                img_path_str = str(img_path)
                if Path(img_path_str).is_file():
                    try:
                        slide.shapes.add_picture(
                            img_path_str,
                            Inches(1),
                            Inches(4.2 + i * 1.5),
                            Inches(4.5),
                        )
                    except Exception as e:
                        log.warning("No se pudo insertar imagen %s: %s", img_path_str, e)

        chart_data = sd.get("chart_data")
        if chart_data:
            _add_chart(
                slide,
                chart_data,
                left=Inches(1),
                top=Inches(3.8),
                width=Inches(8),
                height=Inches(3.2),
            )

        notes_text = str(sd.get("notes", "")).strip()
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text


def read_pptx(file_path: str) -> dict[str, Any]:
    from pptx import Presentation

    path = Path(file_path)
    if not path.is_file():
        return {
            "ok": False,
            "filename": path.name,
            "slide_count": 0,
            "slides": [],
            "error": f"Archivo no encontrado: {file_path}",
        }

    try:
        prs = Presentation(str(path))
        slides_out: list[dict[str, Any]] = []

        for slide in prs.slides:
            title_text = ""
            content_parts: list[str] = []
            images_count = 0

            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    if shape.has_text_frame:
                        title_text = shape.text_frame.text.strip()
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        content_parts.append(text)
                elif shape.shape_type == 13:
                    images_count += 1

            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            slides_out.append({
                "title": title_text,
                "content": "\n".join(content_parts),
                "notes": notes_text,
                "images_count": images_count,
            })

        return {
            "ok": True,
            "filename": path.name,
            "slide_count": len(slides_out),
            "slides": slides_out,
            "error": None,
        }

    except Exception as e:
        log.exception("Error leyendo PPTX: %s", file_path)
        return {
            "ok": False,
            "filename": path.name,
            "slide_count": 0,
            "slides": [],
            "error": f"Error al leer presentación: {e}",
        }


def generate_pptx(
    *,
    title: str,
    slides_data: list[dict[str, Any]],
    template: str = "default",
    folder: str | None = None,
) -> dict[str, Any]:
    from pptx import Presentation

    if not slides_data:
        return {
            "ok": False,
            "filename": None,
            "path": None,
            "size_bytes": 0,
            "slide_count": 0,
            "error": "slides_data está vacío. Se necesita al menos una slide.",
        }

    filepath = resolve_output_path(kind="pptx", title=title, extension="pptx", folder=folder)
    filename = filepath.name

    try:
        prs = Presentation()
        prs.core_properties.title = title
        prs.core_properties.author = "DOT IA"

        _add_title_slide(prs, title)
        _build_default_template(prs, slides_data)

        prs.save(str(filepath))
        size_bytes = filepath.stat().st_size

        log.info(
            "Presentación generada: %s (%d bytes, %d slides + portada)",
            filename,
            size_bytes,
            len(slides_data),
        )

        return {
            "ok": True,
            "filename": filename,
            "path": str(filepath),
            "size_bytes": size_bytes,
            "slide_count": len(slides_data) + 1,
            "error": None,
        }

    except Exception as e:
        log.exception("Error generando PPTX: %s", title)
        return {
            "ok": False,
            "filename": None,
            "path": None,
            "size_bytes": 0,
            "slide_count": 0,
            "error": f"Error al generar presentación: {e}",
        }
